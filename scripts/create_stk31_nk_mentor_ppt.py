from pathlib import Path
import csv
import subprocess

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


ROOT = Path(r"E:\ZHUWEIYU\Documents\R")
OUT_DIR = ROOT / "results" / "STK31_NK_mentor_ppt"
PNG_DIR = OUT_DIR / "png"
PPTX_PATH = OUT_DIR / "STK31_NK_导师汇报_图选择与机制解释.pptx"

PDFTOPPM = Path(r"C:\Users\ZHUWEIYU\.cache\codex-runtimes\codex-primary-runtime\dependencies\native\poppler\Library\bin\pdftoppm.exe")

FOCUSED = ROOT / "results" / "merged_cellchat_focused_stk31_nk"
CELLCHAT = ROOT / "results" / "merged_cellchat_go_plots"
MERGED = ROOT / "results" / "merged_stk31_nk_analysis"

FIGURES = {
    "umap_celltype": MERGED / "umap_manual_celltype_annotation.pdf",
    "umap_stk31": MERGED / "umap_stk31_expression.pdf",
    "umap_groups": MERGED / "umap_stk31_high_and_nk_groups.pdf",
    "focused_circle": FOCUSED / "stk31_nk_mechanism_axis_circle.pdf",
    "focused_bubble": FOCUSED / "stk31_nk_mechanism_axes_bubble.pdf",
    "focused_heatmap": FOCUSED / "stk31_nk_mechanism_axis_heatmap.pdf",
    "go_heatmap": CELLCHAT / "merged_go_enrichment_heatmap.pdf",
}


def ensure_dirs():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    PNG_DIR.mkdir(parents=True, exist_ok=True)


def pdf_to_png(name: str, pdf_path: Path) -> Path:
    if not pdf_path.exists():
        raise FileNotFoundError(f"Missing figure: {pdf_path}")
    out_prefix = PNG_DIR / name
    png_path = PNG_DIR / f"{name}-1.png"
    if png_path.exists() and png_path.stat().st_mtime >= pdf_path.stat().st_mtime:
        return png_path
    subprocess.run(
        [str(PDFTOPPM), "-png", "-singlefile", "-r", "220", str(pdf_path), str(out_prefix)],
        check=True,
    )
    singlefile_png = PNG_DIR / f"{name}.png"
    if singlefile_png.exists():
        singlefile_png.replace(png_path)
    if not png_path.exists():
        raise FileNotFoundError(f"PDF conversion failed: {pdf_path}")
    return png_path


def load_csv(path: Path):
    with path.open("r", encoding="utf-8-sig", newline="") as f:
        return list(csv.DictReader(f))


def load_summary():
    axis_rows = load_csv(FOCUSED / "stk31_high_tumor_epithelial_nk_mechanism_axis_summary.csv")
    pair_rows = load_csv(FOCUSED / "stk31_high_tumor_epithelial_nk_pair_level_cellchat_summary.csv")
    mech_rows = load_csv(FOCUSED / "stk31_high_tumor_epithelial_nk_mechanism_axes_cellchat.csv")
    return axis_rows, pair_rows, mech_rows


def add_textbox(slide, text, left, top, width, height, font_size=18, bold=False, color=(30, 30, 30), align=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, title, Inches(0.45), Inches(0.25), Inches(12.45), Inches(0.5), 24, True, (18, 64, 98))
    if subtitle:
        add_textbox(slide, subtitle, Inches(0.48), Inches(0.78), Inches(12.1), Inches(0.35), 11, False, (90, 90, 90))


def add_bullets(slide, bullets, left, top, width, height, font_size=15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.space_after = Pt(6)
    return box


def add_note(slide, text):
    add_textbox(slide, text, Inches(0.55), Inches(6.72), Inches(12.15), Inches(0.33), 10, False, (110, 110, 110))


def add_image_fit(slide, image_path, left, top, width, height):
    pic = slide.shapes.add_picture(str(image_path), left, top)
    scale = min(width / pic.width, height / pic.height)
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    pic.left = int(left + (width - pic.width) / 2)
    pic.top = int(top + (height - pic.height) / 2)
    return pic


def add_table(slide, rows, headers, left, top, width, height, font_size=11):
    table = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height).table
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(18, 64, 98)
        for p in cell.text_frame.paragraphs:
            p.font.name = "Microsoft YaHei"
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
    for i, row in enumerate(rows, start=1):
        for j, value in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(value)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Microsoft YaHei"
                p.font.size = Pt(font_size)
                p.font.color.rgb = RGBColor(35, 35, 35)
    return table


def fmt_float(value, digits=4):
    return f"{float(value):.{digits}f}"


def build_ppt():
    ensure_dirs()
    pngs = {name: pdf_to_png(name, path) for name, path in FIGURES.items()}
    axis_rows, pair_rows, mech_rows = load_summary()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1. Title
    slide = prs.slides.add_slide(blank)
    add_textbox(slide, "STK31高表达肿瘤上皮细胞与NK细胞互作证据", Inches(0.8), Inches(1.1), Inches(11.8), Inches(0.7), 31, True, (18, 64, 98), PP_ALIGN.CENTER)
    add_textbox(slide, "给导师汇报用：为什么选这几张图，以及每张图说明什么", Inches(1.15), Inches(2.0), Inches(11.0), Inches(0.55), 18, False, (80, 80, 80), PP_ALIGN.CENTER)
    add_bullets(slide, [
        "核心问题：STK31高表达肿瘤上皮细胞是否与NK细胞存在可解释的细胞互作？",
        "分析策略：merged数据作为主分析；先限定肿瘤上皮细胞，再按STK31表达分高/低，CellChat聚焦STK31高表达肿瘤上皮细胞与NK细胞双向通信。",
        "汇报重点：不展示所有结果，只展示能支撑机制假说的图。",
    ], Inches(1.6), Inches(3.0), Inches(10.2), Inches(1.4), 16)
    add_note(slide, "建议开场：这次不是单纯看差异基因，而是围绕 STK31-HLA-I-NK 机制去筛选能解释互作的图。")

    # 2. Storyline
    slide = prs.slides.add_slide(blank)
    add_title(slide, "汇报逻辑：先证明对象，再证明互作，再提出机制")
    add_bullets(slide, [
        "1. 细胞对象：STK31高表达肿瘤上皮细胞和NK细胞在merged数据中可定位。",
        "2. 互作存在：CellChat显示两类细胞存在双向通信，STK31高表达肿瘤上皮细胞到NK方向更强。",
        "3. 机制聚焦：优先解释MHC-I/HLA、TGFb、IFNG/IFN、TIGIT/NECTIN四条轴。",
        "4. 论文呈现：正文放聚焦机制图；全局CellChat和GO图作为补充支持。",
    ], Inches(0.85), Inches(1.45), Inches(5.6), Inches(4.2), 18)
    add_image_fit(slide, pngs["focused_circle"], Inches(6.85), Inches(1.15), Inches(5.75), Inches(4.8))
    add_note(slide, "这一页先把导师带到主线：我们不是为了做很多图，而是选择能回答机制问题的图。")

    # 3. Cell identity and STK31 localization
    slide = prs.slides.add_slide(blank)
    add_title(slide, "图1：先说明细胞定位和STK31表达位置", "为什么要放：导师需要先相信后面CellChat分析中的细胞对象是明确的。")
    add_image_fit(slide, pngs["umap_celltype"], Inches(0.55), Inches(1.15), Inches(4.0), Inches(4.35))
    add_image_fit(slide, pngs["umap_stk31"], Inches(4.65), Inches(1.15), Inches(4.0), Inches(4.35))
    add_image_fit(slide, pngs["umap_groups"], Inches(8.75), Inches(1.15), Inches(4.0), Inches(4.35))
    add_bullets(slide, [
        "讲解点：manual annotation先定义肿瘤上皮细胞，再在肿瘤上皮细胞内部按STK31表达分为高表达和低表达。",
        "放图目的：避免直接跳到CellChat，让老师先看到“分析对象”是清楚的。",
    ], Inches(0.8), Inches(5.75), Inches(11.9), Inches(0.75), 13)

    # 4. Why focused CellChat
    slide = prs.slides.add_slide(blank)
    add_title(slide, "为什么不用全局CellChat图做主图？")
    add_bullets(slide, [
        "全局CellChat图适合说明整体微环境互作，但信息太杂，不能直接回答STK31/NK机制问题。",
        "论文正文应聚焦：STK31_high_tumor_epithelial_to_NK_cell 与 NK_cell_to_STK31_high_tumor_epithelial。",
        "这两个方向的pair-level结果：STK31高表达肿瘤上皮到NK、NK到STK31高表达肿瘤上皮。",
        "因此正文主图应放聚焦版circle/bubble/heatmap，全局CellChat放补充材料。",
    ], Inches(0.85), Inches(1.25), Inches(6.0), Inches(4.8), 17)
    pair_table_rows = []
    for row in pair_rows:
        pair_table_rows.append([row["metric"], row["direction"], fmt_float(row["value"], 4)])
    add_table(slide, pair_table_rows, ["指标", "方向", "值"], Inches(7.05), Inches(1.45), Inches(5.25), Inches(1.7), 9)
    add_image_fit(slide, pngs["focused_circle"], Inches(7.05), Inches(3.35), Inches(5.25), Inches(2.75))
    add_note(slide, "讲解时强调：全局图是背景证据，聚焦图才是正文机制证据。")

    # 5. Circle figure
    slide = prs.slides.add_slide(blank)
    add_title(slide, "图2：机制轴circle图用于概览双向互作")
    add_image_fit(slide, pngs["focused_circle"], Inches(0.6), Inches(1.05), Inches(6.2), Inches(5.35))
    add_bullets(slide, [
        "为什么放：一眼展示STK31高表达肿瘤上皮细胞与NK细胞存在双向通信。",
        "它回答的问题：互作方向是什么？哪些机制轴参与？",
        "主要结论：MHC-I/HLA轴以STK31高表达肿瘤上皮到NK为主；TGFb、IFNG/IFN体现NK到肿瘤上皮的反馈；TIGIT/NECTIN双向可见。",
        "定位：适合做正文概览图，但具体配体受体证据需要bubble图支撑。",
    ], Inches(7.1), Inches(1.35), Inches(5.45), Inches(4.4), 16)

    # 6. Bubble figure
    slide = prs.slides.add_slide(blank)
    add_title(slide, "图3：bubble图是最核心的机制证据")
    add_image_fit(slide, pngs["focused_bubble"], Inches(0.45), Inches(1.05), Inches(7.2), Inches(5.6))
    add_bullets(slide, [
        "为什么放：它展示具体ligand-receptor，而不是只给总强度。",
        "MHC-I/HLA：HLA-A/B/C/E/F -> CD8A 或 KLRK1，支持HLA-I相关互作。",
        "TGFb：TGFB1 -> TGFb receptor复合体，提示免疫调节/抑制反馈。",
        "IFNG：IFNG -> IFNGR1/IFNGR2，提示NK细胞对肿瘤上皮的炎症反馈。",
        "TIGIT/NECTIN：NECTIN2/3 -> TIGIT 和 TIGIT -> NECTIN2，支持检查点互作。",
    ], Inches(8.0), Inches(1.25), Inches(4.9), Inches(5.15), 14)

    # 7. Axis summary
    slide = prs.slides.add_slide(blank)
    add_title(slide, "图4：heatmap/summary说明哪条机制轴最强")
    add_image_fit(slide, pngs["focused_heatmap"], Inches(0.65), Inches(1.05), Inches(5.4), Inches(4.7))
    axis_table_rows = [[r["mechanism_axis"], r["direction"], fmt_float(r["total_probability"], 4), r["n_interactions"]] for r in axis_rows]
    add_table(slide, axis_table_rows, ["机制轴", "方向", "总概率", "互作数"], Inches(6.25), Inches(1.25), Inches(6.55), Inches(2.75), 8)
    add_bullets(slide, [
        "结论1：MHC-I/HLA轴是最强的STK31高表达肿瘤上皮到NK方向信号。",
        "结论2：TGFb和IFNG/IFN更偏向NK到肿瘤上皮反馈。",
        "结论3：TIGIT/NECTIN是免疫检查点方向，适合做后续验证。",
    ], Inches(6.35), Inches(4.25), Inches(6.25), Inches(1.55), 14)
    add_note(slide, "这页适合回答导师：为什么你选这四条机制轴，而不是随便挑通路。")

    # 8. GO support
    slide = prs.slides.add_slide(blank)
    add_title(slide, "GO图的定位：作为功能背景支持，不作为互作主证据")
    add_image_fit(slide, pngs["go_heatmap"], Inches(0.7), Inches(1.15), Inches(6.0), Inches(4.8))
    add_bullets(slide, [
        "为什么可以放：GO可以补充说明差异基因背后的生物学状态。",
        "为什么不作为主图：GO不能直接证明两类细胞存在配体-受体互作。",
        "建议放法：若GO结果能突出免疫调节、抗原呈递、IFN反应，可放正文末尾或补充图。",
        "和CellChat关系：GO解释功能背景，CellChat解释细胞间通信。",
    ], Inches(7.1), Inches(1.35), Inches(5.45), Inches(4.2), 16)

    # 9. Suggested final figure set
    slide = prs.slides.add_slide(blank)
    add_title(slide, "建议给导师看的最终图组合")
    add_bullets(slide, [
        "正文主图A：UMAP/FeaturePlot，说明STK31高/低表达肿瘤上皮细胞和NK细胞的空间定位。",
        "正文主图B：机制轴circle，快速说明两类细胞存在双向通信。",
        "正文主图C：机制轴bubble，展示MHC-I/HLA、TGFb、IFNG、TIGIT/NECTIN的具体配体受体。",
        "正文主图D：机制轴heatmap或summary，说明各机制轴方向性和强弱。",
        "补充图：全局CellChat circle/heatmap、GO heatmap/dotplot、完整通信表。",
    ], Inches(0.85), Inches(1.2), Inches(6.05), Inches(4.9), 16)
    add_image_fit(slide, pngs["focused_bubble"], Inches(7.1), Inches(1.05), Inches(5.55), Inches(4.4))
    add_note(slide, "导师沟通重点：正文图少而准，补充图完整但不抢主线。")

    # 10. Next validation
    slide = prs.slides.add_slide(blank)
    add_title(slide, "下一步验证建议：从计算证据走向实验闭环")
    add_bullets(slide, [
        "1. 细胞模型：胆囊癌细胞中敲低/过表达STK31。",
        "2. HLA-I验证：检测HLA-A/B/C/E表达，建议流式 + qPCR/Western。",
        "3. NK共培养：检测NK杀伤、CD107a脱颗粒、IFNG释放。",
        "4. 机制阻断/救援：重点看TIGIT-NECTIN2阻断、IFNG/HLA-I变化、TGFb相关抑制信号。",
        "5. 论文表述：CellChat提示互作机制，实验验证STK31是否调控HLA-I并影响NK功能。",
    ], Inches(0.95), Inches(1.2), Inches(11.6), Inches(4.9), 18)
    add_note(slide, "这页用于结尾：目前是机制线索，下一步实验验证可以把故事补成因果链。")

    prs.save(PPTX_PATH)
    return PPTX_PATH


if __name__ == "__main__":
    print(build_ppt())
